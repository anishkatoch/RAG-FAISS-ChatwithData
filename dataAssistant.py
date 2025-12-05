import streamlit as st
import faiss
import os
from io import BytesIO
from docx import Document
import numpy as np
from PyPDF2 import PdfReader
from langchain.chains import RetrievalQA
from langchain.memory import ConversationSummaryBufferMemory
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.embeddings import OpenAIEmbeddings
from langchain_community.vectorstores import FAISS
from langchain_community.docstore.in_memory import InMemoryDocstore
from langchain_community.chat_models import ChatOpenAI
from dotenv import load_dotenv
from pptx import Presentation
# Load environment variables
load_dotenv()
openai_api_key = os.getenv("OPENAI_API_KEY")
if not openai_api_key:
    st.error("OpenAI API Key is not set. Please add it to your environment variables.")
    st.stop()

MAX_MESSAGES = 5  # Maximum number of messages to retain


def answer_question(vectorstore, query, memory):
    """
    Answers a question based on the vectorstore and the conversation memory.
    Args:
        vectorstore: A FAISS vector store instance for retrieving context.
        query: The user's question/query.
        memory: A memory object storing the chat history.
    Returns:
        A string response from the model.
    """
    try:
        # Initialize the language model
        llm = ChatOpenAI(model="gpt-4", openai_api_key=openai_api_key, temperature=0.2, max_tokens=150)

        # Setup the retriever using the vector store
        retriever = vectorstore.as_retriever(search_kwargs={"k": 2})  # Retrieve top 2 documents

        # Build the RetrievalQA chain with memory integration
        qa_chain = RetrievalQA.from_chain_type(
            llm=llm,
            retriever=retriever,
            memory=memory,  # Provides conversation context
            chain_type="stuff",  # Recommended default chain type
            return_source_documents=False,  # Turn off raw documents return (optional)
        )

        # Run the chain and return the output
        response = qa_chain.invoke({"query": query})
        return response["result"]

    except ValueError as ve:
        raise ValueError(f"ValueError in answer_question: {ve}")
    except AttributeError as ae:
        raise AttributeError(f"AttributeError in answer_question: {ae}")
    except Exception as e:
        raise RuntimeError(f"Unexpected error in answer_question: {e}")



def process_input(input_data):
    """Processes different input types and returns a vectorstore."""
    texts = []
    for data in input_data:
        input_type, content = data

        if input_type == "PDF":
            pdf_reader = PdfReader(BytesIO(content.read()))
            for page in pdf_reader.pages:
                page_text = page.extract_text()
                if page_text:
                    texts.append(page_text)
        elif input_type == "DOCX":
            doc = Document(BytesIO(content.read()))
            for para in doc.paragraphs:
                if para.text:
                    texts.append(para.text)
        elif input_type == "TXT":
            txt_content = content.read().decode("utf-8")
            texts.append(txt_content)
        elif input_type == "PPTX":
            ppt = Presentation(BytesIO(content.read()))
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        texts.append(shape.text)

    if not texts:
        raise ValueError("No text content extracted from the provided files.")

    # Use RecursiveCharacterTextSplitter for efficient splitting
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
    split_texts = text_splitter.split_text(" ".join(texts))

    embeddings = OpenAIEmbeddings(openai_api_key=openai_api_key)

    # Initialize FAISS vector store
    sample_embedding = np.array(embeddings.embed_query("sample text"))
    dimension = sample_embedding.shape[0]
    index = faiss.IndexFlatL2(dimension)

    # Create a vectorstore with batch processing
    vector_store = FAISS(
        embedding_function=embeddings.embed_query,
        index=index,
        docstore=InMemoryDocstore(),
        index_to_docstore_id={},
    )

    # Batch embed texts to avoid memory overload
    batch_size = 50
    for i in range(0, len(split_texts), batch_size):
        batch = split_texts[i:i + batch_size]
        vector_store.add_texts(batch)

    return vector_store

def truncate_messages(messages, max_messages=5):
    """
    Truncates the message list to retain only the last `max_messages` messages.
    If there are too many messages, summarizes the earlier ones.
    """
    if len(messages) > max_messages:
        # Summarize earlier messages (if needed)
        summary = " ".join(msg["content"] for msg in messages[:-max_messages])
        summarized_message = {
            "role": "assistant",
            "content": f"Summary of previous messages: {summary}"
        }
        # Keep only the summary and last few messages
        return [summarized_message] + messages[-max_messages:]
    return messages


# Modify the main function to handle PPTX uploads
# Add this CSS styling at the beginning of your main function
def main():
    # Custom CSS for styling
    st.markdown("""
        <style>
            .main {
                background-color: #f0f2f6;
            }
            .stChatInput input {
                border-radius: 20px !important;
                padding: 10px !important;
            }
            .stButton>button {
                background-color: #4CAF50;
                color: white;
                border-radius: 8px;
                padding: 10px 24px;
            }
            .stButton>button:hover {
                background-color: #45a049;
            }
            .sidebar .sidebar-content {
                background-color: #2c3e50;
                color: white;
            }
            .stMarkdown h1 {
                color: #2c3e50;
                text-align: center;
                font-size: 2.5em;
            }
            .chat-message {
                padding: 1.5rem;
                border-radius: 0.5rem;
                margin-bottom: 1rem;
                display: flex;
            }
            .chat-message.user {
                background-color: #e3f2fd;
            }
            .chat-message.assistant {
                background-color: #fff3e0;
            }
        </style>
    """, unsafe_allow_html=True)

    # Main app header
    st.markdown(
        """<div style='background: linear-gradient(45deg, #4CAF50, #2196F3); 
        padding: 2rem; border-radius: 1rem; box-shadow: 0 4px 6px rgba(0,0,0,0.1)'>
        <h1 style='color: white; text-align: center; margin: 0;'>📚 Smart Data Assistant</h1>
        <p style='color: white; text-align: center; margin: 0.5rem 0;'>“Retrieval-Augmented Chat with GPT-4, FAISS & LangChain”</p>
        </div>""",
        unsafe_allow_html=True
    )
    # Short description for all users
    st.markdown(
        """
        <div style='margin-top:1em;margin-bottom:1em;padding:1em;background:#e3f2fd;border-radius:8px;'>
        <b>What is this?</b> This app lets anyone chat with the contents of a document using AI. By default, it loads "Sample-AI Overview.docx" so you can ask questions and get instant answers from the document. It’s useful for anyone who wants quick insights or information from documents without reading them fully.
        </div>
        """,
        unsafe_allow_html=True
    )

    with st.sidebar:
        st.markdown("## 🗂️ Document Upload Zone")
        st.markdown("Upload your files here to get started!")

        with st.expander("📤 Upload Documents", expanded=True):
            pdf_files = st.file_uploader("PDF Files", type="pdf",
                                         accept_multiple_files=True,
                                         help="Upload PDF documents")
            docx_files = st.file_uploader("Word Documents", type="docx",
                                          accept_multiple_files=True,
                                          help="Upload Word documents")

        # If user uploads files, show Process Documents button
        files_uploaded = (pdf_files and len(pdf_files) > 0) or (docx_files and len(docx_files) > 0)
        if files_uploaded:
            if st.button("✨ Process Documents", use_container_width=True):
                try:
                    input_data = []
                    if pdf_files:
                        input_data.extend([("PDF", file) for file in pdf_files])
                    if docx_files:
                        input_data.extend([("DOCX", file) for file in docx_files])

                    vectorstore = process_input(input_data)
                    st.session_state["vectorstore"] = vectorstore
                    st.session_state["memory"] = ConversationSummaryBufferMemory(
                        llm=ChatOpenAI(model="gpt-4", openai_api_key=openai_api_key),
                        memory_key="chat_history",
                        max_token_limit=500
                    )
                    st.session_state["messages"] = [
                        {"role": "assistant", "content": "Hello! How can I assist you today?"}
                    ]
                    st.success("Inputs processed successfully!")
                except Exception as e:
                    st.error(f"Error processing input: {e}")
        # If no files uploaded and no vectorstore, auto-process default file
        elif "vectorstore" not in st.session_state:
            try:
                with open('Sample-AI Overview.docx', 'rb') as f:
                    default_file = BytesIO(f.read())
                input_data = [("DOCX", default_file)]
                vectorstore = process_input(input_data)
                st.session_state["vectorstore"] = vectorstore
                st.session_state["memory"] = ConversationSummaryBufferMemory(
                    llm=ChatOpenAI(model="gpt-4", openai_api_key=openai_api_key),
                    memory_key="chat_history",
                    max_token_limit=500
                )
                st.session_state["messages"] = [
                    {"role": "assistant", "content": "Hello! How can I assist you today?"}
                ]
                st.info("No files uploaded. Using default document for chat.")
            except Exception as e:
                st.error(f"Could not read default document: {e}")
    # ... [keep your existing processing code here] ...

    # Chat interface
    if "vectorstore" in st.session_state and "memory" in st.session_state:
        st.markdown("### 💬 Chat with Your Data")
        st.markdown("Ask anything about your uploaded data!")

        # Show sample questions only if default doc is loaded (no user upload this session)
        files_uploaded = False
        with st.sidebar:
            pdf_files = st.session_state.get('pdf_files', None)
            docx_files = st.session_state.get('docx_files', None)
            if (pdf_files and len(pdf_files) > 0) or (docx_files and len(docx_files) > 0):
                files_uploaded = True

        if not files_uploaded:
            st.markdown("#### Example Questions:")
            sample_questions = [
                "What are the three types of AI mentioned in the document?",
                "How is AI used in the finance industry?",
                "For what purposes does Amazon Alexa employ AI?"
            ]
            for q in sample_questions:
                if st.button(q, key=q):
                    st.session_state.messages.append({"role": "user", "content": q})
                    st.chat_message("user").write(q)
                    try:
                        vectorstore = st.session_state["vectorstore"]
                        memory = st.session_state["memory"]
                        response = answer_question(vectorstore, q, memory)
                        st.session_state.messages.append({"role": "assistant", "content": response})
                        st.chat_message("assistant").markdown(response)
                    except Exception as e:
                        st.error(f"Error generating response: {e}")

        # Display chat messages
        for msg in st.session_state.messages:
            avatar = "🤖" if msg["role"] == "assistant" else "👤"
            with st.chat_message(msg["role"], avatar=avatar):
                st.markdown(f"<div class='chat-message {msg['role']}'>{msg['content']}</div>",
                            unsafe_allow_html=True)

        # Chat input
        if prompt := st.chat_input("Ask your question here..."):
            st.session_state.messages.append({"role": "user", "content": prompt})
            st.chat_message("user").write(prompt)

            # Truncate messages after adding a new one
            st.session_state.messages = truncate_messages(st.session_state.messages, MAX_MESSAGES)

            try:
                vectorstore = st.session_state["vectorstore"]
                memory = st.session_state["memory"]
                response = answer_question(vectorstore, prompt, memory)
                st.session_state.messages.append({"role": "assistant", "content": response})
                st.chat_message("assistant").write(response)
            except Exception as e:
                if "tokens" in str(e):
                    st.error("The input or output tokens exceeded the limit. Try shortening your query or context.")
                else:
                    st.error(f"Error generating response: {e}")
    # ... [keep your existing chat handling code here] ...


if __name__ == "__main__":
    main()

